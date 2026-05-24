"""
pptx_parser.py — Phase 1 of the AI Quiz Generator.

Single responsibility: given a PowerPoint file, return all of its text.
No AI here — just reading.
"""

from pptx import Presentation


def extract_text_from_pptx(file) -> str:
    """Pull all the text out of a PowerPoint file and return it as one big string.

    Args:
        file: Either a file path (str) or a Streamlit UploadedFile object.
              python-pptx accepts both, which is why we don't care.

    Returns:
        A single string with all slide text, separated by "--- Slide N ---" markers.
    """
    prs = Presentation(file)
    all_text = []

    for i, slide in enumerate(prs.slides, start=1):
        slide_text = [f"--- Slide {i} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_text.append(text)
        all_text.append("\n".join(slide_text))

    return "\n\n".join(all_text)


if __name__ == "__main__":
    # This block only runs when you do `python pptx_parser.py somefile.pptx`
    # It does NOT run when app.py imports this file. Great for quick testing.
    import sys

    if len(sys.argv) < 2:
        print("Usage: python pptx_parser.py <path-to-pptx>")
        sys.exit(1)

    text = extract_text_from_pptx(sys.argv[1])
    print(text)
